"""Phase G.4 Task 21 — real-process MarkItDown MCP verification.

Drives the disabled production `official-markitdown` catalog entry through the
full G.3/G.4 pipeline under an isolated temp directory:
  1. Create an isolated venv.
  2. Install from the hash-locked dependency file with the exact production
     pip flags (`--require-hashes --no-deps --no-input --disable-pip-version-check`).
  3. Use the installed venv libraries to create real PPTX/XLSX fixtures.
  4. Build a one-off catalog with the entry enabled.
  5. Detect a document-to-markdown request and capture document snapshots.
  6. Prepare a hash-bound auto-provisioning plan.
  7. Auto-approve -> install/validate/activate.
  8. Resume and exercise the exact-file invocation policy through the real
     `mcp.markitdown.convert_to_markdown` tool for every advertised extension.

Reports which catalog extensions actually convert in the reviewed
CPython 3.13 / Windows win_amd64 environment.  Any extension that fails is
reported as MUST_BE_REMOVED from `config/mcp_catalog.json` before the entry
can be enabled.

Isolated under a temp base_dir/managed_root — never touches the real
app_data/mcp_servers/ state or config/mcp_catalog.json.

Run: venv/Scripts/python.exe scripts/manual_verify_g4_markitdown.py
"""
import copy
import json
import os
import shutil
import subprocess
import sys
import tempfile

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from mcp_layer.runtime_manager import MultiMcpRuntimeManager, RuntimeState  # noqa: E402
from mcp_management.auto_provisioning import AutoProvisioningManager  # noqa: E402
from mcp_management.catalog import build_catalog  # noqa: E402
from mcp_management.provisioning_models import AutoProvisioningApproval  # noqa: E402
from mcp_management.registry import get_installed  # noqa: E402
from tools.executor import ToolExecutor  # noqa: E402
from tools.models import ToolCall  # noqa: E402
from tools.registry import default_registry  # noqa: E402


_CATALOG_PATH = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                             "config", "mcp_catalog.json")
_REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

_EXTENSIONS_TO_TEST = (".txt", ".html", ".htm", ".pdf", ".docx", ".pptx", ".xls", ".xlsx")


def _ok(label):
    print(f"[OK] {label}")


def _warn(label):
    print(f"[WARN] {label}")


def _fail(label):
    print(f"[FAIL] {label}")


def _load_production_markitdown_entry():
    with open(_CATALOG_PATH, "r", encoding="utf-8") as f:
        data = json.load(f)
    entry = data["servers"]["official-markitdown"]
    if entry.get("enabled"):
        raise RuntimeError(
            "The production catalog entry is already enabled.  Run this script only while it is disabled.")
    return entry


def _build_enabled_catalog(entry_raw):
    """Return a catalog with a temporary enabled copy of the entry."""
    enabled = copy.deepcopy(entry_raw)
    enabled["enabled"] = True
    return build_catalog({"catalog_version": 1, "servers": {"official-markitdown-enabled": enabled}})


def _create_venv(venv_dir: str) -> str:
    """Create a CPython venv and return the python executable path."""
    subprocess.run([sys.executable, "-m", "venv", venv_dir], check=True)
    if sys.platform == "win32":
        return os.path.join(venv_dir, "Scripts", "python.exe")
    return os.path.join(venv_dir, "bin", "python")


def _install_lock(venv_python: str, lock_path: str) -> None:
    """Install the curated lock with the exact production flags."""
    cmd = [
        venv_python, "-m", "pip", "install",
        "--require-hashes", "--no-deps", "--no-input", "--disable-pip-version-check",
        "-r", lock_path,
    ]
    print(f"Running: {' '.join(cmd)}")
    subprocess.run(cmd, check=True)


def _create_fixtures_with_venv(venv_python: str, fixtures_dir: str) -> dict[str, bool]:
    """Populate the isolated fixtures dir, preferring committed fixtures.

    Returns a mapping of extension -> whether a fixture is present at the end.
    """
    os.makedirs(fixtures_dir, exist_ok=True)
    created = {}

    # Phase G.4: copy every committed markitdown_sample.* fixture first so the
    # manual verification tests the actual files that would be used in CI.
    committed_src = os.path.join(_REPO_ROOT, "tests", "fixtures")
    for name in os.listdir(committed_src):
        if name.startswith("markitdown_sample."):
            src = os.path.join(committed_src, name)
            if os.path.isfile(src):
                shutil.copy2(src, fixtures_dir)

    # Any missing formats are generated using only libraries available in the
    # reviewed venv, so the script remains self-contained for new formats.

    # TXT and HTML need no external libraries.
    for ext, content in (
        (".txt", "G4-VERIFY-TXT-2026 MarkItDown text fixture.\n"),
        (".html", (
            "<!DOCTYPE html><html><head><title>Fixture</title></head>"
            "<body><h1>G4-VERIFY-HTML-2026</h1><p>MarkItDown HTML fixture.</p></body></html>"
        )),
    ):
        path = os.path.join(fixtures_dir, f"markitdown_sample{ext}")
        if not os.path.isfile(path):
            with open(path, "w", encoding="utf-8") as f:
                f.write(content)
            _ok(f"created {ext} fixture")

    # PPTX and XLSX using venv libraries.
    for ext, script in (
        (".pptx", """
import sys, os
from pptx import Presentation
from pptx.util import Inches
out = sys.argv[1]
os.makedirs(os.path.dirname(out), exist_ok=True)
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
left = top = Inches(1)
box = slide.shapes.add_textbox(left, top, Inches(4), Inches(1))
box.text_frame.text = "G4-VERIFY-PPTX-2026 MarkItDown PowerPoint fixture"
prs.save(out)
"""),
        (".xlsx", """
import sys, os
import openpyxl
out = sys.argv[1]
os.makedirs(os.path.dirname(out), exist_ok=True)
wb = openpyxl.Workbook()
ws = wb.active
ws["A1"] = "G4-VERIFY-XLSX-2026"
ws["A2"] = "MarkItDown Excel fixture"
wb.save(out)
"""),
    ):
        path = os.path.join(fixtures_dir, f"markitdown_sample{ext}")
        if os.path.isfile(path):
            continue
        try:
            subprocess.run([venv_python, "-c", script, path], check=True,
                          capture_output=True, text=True)
            _ok(f"created {ext} fixture")
        except subprocess.CalledProcessError as exc:
            _warn(f"could not create {ext} fixture: {exc.stderr.strip() or exc}")

    # DOCX: python-docx is intentionally omitted from the curated lock.
    # Build a minimal package and see if markitdown can still convert it.
    docx_path = os.path.join(fixtures_dir, "markitdown_sample.docx")
    if not os.path.isfile(docx_path):
        if _make_minimal_docx(docx_path):
            _ok("created .docx fixture (minimal package)")
        else:
            _warn("could not create .docx fixture")

    # XLS: xlwt is not in the curated lock; a valid binary .xls cannot be
    # generated without it.  Any pre-existing fixture would need to be copied.
    xls_src = os.path.join(_REPO_ROOT, "tests", "fixtures", "markitdown_sample.xls")
    xls_dst = os.path.join(fixtures_dir, "markitdown_sample.xls")
    if not os.path.isfile(xls_dst) and os.path.isfile(xls_src):
        shutil.copy2(xls_src, xls_dst)
        _ok("copied pre-existing .xls fixture")

    # Final status for every extension we intend to test.
    for ext in _EXTENSIONS_TO_TEST:
        created[ext] = os.path.isfile(os.path.join(fixtures_dir, f"markitdown_sample{ext}"))

    return created


def _make_minimal_docx(path: str) -> bool:
    import zipfile

    parts = {
        "[Content_Types].xml": (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
            '<Default Extension="xml" ContentType="application/xml"/>'
            '<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
            '</Types>'
        ),
        "_rels/.rels": (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>'
            '</Relationships>'
        ),
        "word/_rels/document.xml.rels": (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '</Relationships>'
        ),
        "word/document.xml": (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
            '<w:body><w:p><w:r><w:t>G4-VERIFY-DOCX-2026 MarkItDown DOCX fixture</w:t></w:r></w:p></w:body>'
            '</w:document>'
        ),
    }
    try:
        with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as zf:
            for name, data in parts.items():
                zf.writestr(name, data.encode("utf-8"))
        return True
    except Exception as exc:
        _warn(f"could not build minimal DOCX: {exc}")
        return False


def _file_uri(local_path: str) -> str:
    import urllib.parse
    import urllib.request

    return urllib.parse.urljoin("file:", urllib.request.pathname2url(os.path.abspath(local_path)))


def _test_extension(runtime_manager, reg, fixtures_dir: str, ext: str) -> tuple[bool, str]:
    """Exercise the installed server against one fixture extension."""
    fixture_path = os.path.join(fixtures_dir, f"markitdown_sample{ext}")
    if not os.path.isfile(fixture_path):
        return False, "fixture missing"

    marker = f"G4-VERIFY-{ext.lstrip('.').upper()}-2026"
    from mcp_management.document_authorization import DocumentAuthorizationStore, DocumentInputSnapshot
    from tools.models import ToolPermission

    stat = os.stat(fixture_path)
    snap = DocumentInputSnapshot(
        source_uri=fixture_path,
        local_path=fixture_path,
        size_bytes=stat.st_size,
        ctime_ns=int(stat.st_ctime_ns),
        permission=ToolPermission.READ,
    )
    auth = DocumentAuthorizationStore.default().create_authorization(snap)

    try:
        executor = ToolExecutor(reg)
        file_uri = _file_uri(fixture_path)
        call = ToolCall(call_id=f"g4_{ext.lstrip('.')}",
                        tool_name="mcp.markitdown.convert_to_markdown",
                        arguments={"uri": file_uri})
        result = executor.execute(call)
        if not result.success:
            return False, f"tool failure ({result.error_code}): {result.message}"
        text = result.data.get("text", "") if isinstance(result.data, dict) else ""
        if marker not in text:
            return False, f"marker {marker!r} not found in conversion output"
        return True, "converted and marker found"
    except Exception as exc:
        return False, f"exception: {exc}"
    finally:
        try:
            DocumentAuthorizationStore.default().consume_authorization(auth.auth_id)
        except Exception:
            pass


def main():
    entry_raw = _load_production_markitdown_entry()
    catalog = _build_enabled_catalog(entry_raw)
    entry = catalog.get("official-markitdown-enabled")
    reg = default_registry()

    tmp_root = tempfile.mkdtemp(prefix="g4_markitdown_manual_")
    try:
        managed_root = "app_data/mcp_servers"
        fixtures_dir = os.path.join(tmp_root, "tests", "fixtures")

        # Phase 0: isolated venv + curated lock install.
        venv_dir = os.path.join(tmp_root, ".venv")
        venv_python = _create_venv(venv_dir)
        lock_path = os.path.join(_REPO_ROOT, entry_raw["installer"]["lock_file"])
        _install_lock(venv_python, lock_path)
        _ok(f"installed lock file into isolated venv: {venv_dir}")

        # Phase 1: create fixtures using only venv libraries.
        fixture_status = _create_fixtures_with_venv(venv_python, fixtures_dir)

        # Phase 2: provision with project code (uses the same venv via launch spec).
        manager = AutoProvisioningManager(catalog, base_dir=tmp_root, managed_root=managed_root,
                                          registry_path=None)
        runtime_manager = MultiMcpRuntimeManager(reg, base_dir=tmp_root, managed_root=managed_root)
        _ok("catalog + registries isolated")

        user_text = f"convert {fixtures_dir}{os.sep}markitdown_sample.pdf to markdown"
        from mcp_management.document_authorization import build_document_snapshots_from_text
        document_snapshots = build_document_snapshots_from_text(user_text)
        request = manager.begin_request(user_text, "document_to_markdown", entry,
                                        document_snapshots=document_snapshots)
        assert request is not None, "begin_request returned None (entry not eligible)"
        assert len(request.document_snapshots) == 1, "expected one PDF snapshot captured"
        plan = manager.prepare_plan(request.request_id)
        assert plan.lock_file_hash is not None, "lock file hash must be computed"
        print("\n".join(plan.summary_lines()))
        _ok("provisioning plan prepared with document snapshot bound into hash")

        approval = AutoProvisioningApproval(approved=True, plan_id=plan.plan_id,
                                          plan_hash=plan.compute_hash())
        result = manager.provision_and_activate(request.request_id, runtime_manager,
                                                 approval=approval)
        assert result.installed_version == entry_raw["installer"]["version"]
        installed = get_installed(entry.server_id, None, tmp_root, managed_root)
        assert installed is not None and installed.installer_type == "python_venv"
        status = runtime_manager.get_status(entry.server_id)
        assert status.state == RuntimeState.HEALTHY
        _ok(f"installed + validated + activated; runtime {status.state.value}, "
           f"{status.registered_tool_count} tool(s) registered")

        resumed_text = manager.resume(request.request_id)
        assert resumed_text == user_text
        _ok("original request resumed; PDF authorization created")

        # Phase 3: exercise every advertised extension through the real tool.
        print("\n--- extension conversion results ---")
        supported = []
        unsupported = []
        for ext in _EXTENSIONS_TO_TEST:
            if not fixture_status.get(ext):
                _fail(f"{ext}: fixture unavailable")
                unsupported.append(ext)
                continue
            ok, detail = _test_extension(runtime_manager, reg, fixtures_dir, ext)
            if ok:
                _ok(f"{ext}: {detail}")
                supported.append(ext)
            else:
                _fail(f"{ext}: {detail}")
                unsupported.append(ext)

        # Phase 4: report.
        print("\n=== G.4 MarkItDown manual verification report ===")
        print(f"Temporary base dir:    {tmp_root}")
        print(f"Venv Python:           {venv_python}")
        print(f"Installed version:     {result.installed_version}")
        print(f"Candidate validator:   {entry_raw.get('candidate_validator')}")
        print(f"Invocation policy:     {entry_raw.get('invocation_policy')}")
        print(f"Supported extensions   ({len(supported)}): {supported}")
        print(f"Unsupported extensions ({len(unsupported)}): {unsupported}")
        if unsupported:
            print("\nACTION REQUIRED: remove these extensions from the production catalog entry's")
            print("`selection_hints.extensions.document_to_markdown` list before enabling the entry:")
            for ext in unsupported:
                print(f"  - {ext}")
            sys.exit(1)
        print("\nAll advertised extensions converted successfully. The catalog entry may be enabled.")

    finally:
        try:
            runtime_manager.stop_all()
        except Exception:
            pass
        shutil.rmtree(tmp_root, ignore_errors=True)


if __name__ == "__main__":
    main()
